"""
Extrae la imagen principal (retrato) de cada slide del PPTX
y la guarda en 50_visionarios/images/slide_XX.jpg

Lógica:
- En cada slide buscamos la imagen más grande (el retrato)
- La guardamos como slide_01.jpg, slide_02.jpg, etc.
- Al final imprime un mapeo slide -> nombre de persona
"""
import os
import sys
import zipfile
import shutil
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import io

if sys.stdout.encoding != 'utf-8':
    sys.stdout.reconfigure(encoding='utf-8')

PPT_PATH = "ppt.pptx"
OUTPUT_DIR = os.path.join("50_visionarios", "images")

# Mapeo slide_index (0-based) -> nombre de persona (para el log)
NAMES = [
    "Leopoldo Izquierdo",        # 0
    "Jose Joaquin",               # 1
    "Irene Maclino Navarro",      # 2
    "Gonzalo Herreros Moya",      # 3
    "Sebastian de la Obra",       # 4
    "Antonio Vallejo Triano",     # 5
    "Miguel Roldan",              # 6
    "Manuel Murillo Estevez",     # 7
    "Antonio Monterroso Checa",   # 8
    "Isabel Albas Vives",         # 9
    "Marian Aguilar",             # 10
    "Eduardo Lucena",             # 11
    "Anselmo Cordoba",            # 12
    "Catalina Molina Rodriguez",  # 13
    "Antonio Cano",               # 14
    "Elena Rizos Espinosa",       # 15
    "Antonia Alcantara",          # 16
    "Paco Morales",               # 17
    "Leonardo Gallardo Apolo",    # 18
    "Antonio Munoz Ariza",        # 19
    "Rafael Carrillo",            # 20
    "Paco Rosales",               # 21
    "Javier Campos",              # 22
    "Rafael San Miguel",          # 23
    "Ricardo Rojas Peinado",      # 24
    "Almudena Villegas",          # 25
    "Javier Alcala de la Moneda", # 26
    "Rafael Muela Rodriguez",     # 27
    "Macarena Sanchez del Aguila",# 28
    "Teresa Guitar",              # 29
    "David Pino",                 # 30
    "Lola Perez",                 # 31
    "Antonio Manuel Rodriguez",   # 32
    "Desiderio Vaquerizo Gil",    # 33
    "Ma Angeles Jordano Barbudo", # 34
    "Ricardo Hernandez",          # 35
    "Carmen Balbuena",            # 36
    "Manuel Rivera Mateos",       # 37
    "Nuria Ceular Villamandos",   # 38
    "Jose Antonio Fernandez Gallardo", # 39
    "Teresa Avalos Urena",        # 40
    "Paloma Lopez-Sidro",         # 41
    "Walada de la Mata",          # 42
    # extras si hay mas slides
]

def main():
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    
    if not os.path.exists(PPT_PATH):
        print(f"ERROR: No se encuentra {PPT_PATH}")
        return
    
    print(f"Abriendo {PPT_PATH}...")
    prs = Presentation(PPT_PATH)
    print(f"Total de slides: {len(prs.slides)}")
    
    results = []
    
    for slide_idx, slide in enumerate(prs.slides):
        name = NAMES[slide_idx] if slide_idx < len(NAMES) else f"Slide_{slide_idx}"
        
        # Recopilar todas las imágenes del slide con su tamaño
        images_in_slide = []
        for shape in slide.shapes:
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                try:
                    img_blob = shape.image.blob
                    img_format = shape.image.ext  # jpg, png, etc.
                    # Tamaño en EMU (mayor = imagen más grande)
                    area = shape.width * shape.height
                    images_in_slide.append((area, img_blob, img_format, shape))
                except Exception as e:
                    pass
        
        if not images_in_slide:
            print(f"  Slide {slide_idx+1:02d} [{name}]: SIN IMAGEN")
            results.append((slide_idx, name, None))
            continue
        
        # Tomar la imagen más grande (el retrato)
        images_in_slide.sort(key=lambda x: x[0], reverse=True)
        _, img_blob, img_format, _ = images_in_slide[0]
        
        # Guardar la imagen
        out_filename = f"slide_{slide_idx+1:02d}.jpg"
        out_path = os.path.join(OUTPUT_DIR, out_filename)
        
        try:
            # Convertir a JPEG si es PNG/otro formato
            img = Image.open(io.BytesIO(img_blob))
            # Convertir a RGB si es RGBA (para JPEG)
            if img.mode in ('RGBA', 'P', 'LA'):
                img = img.convert('RGB')
            img.save(out_path, 'JPEG', quality=90)
            
            w, h = img.size
            print(f"  Slide {slide_idx+1:02d} [{name}]: {out_filename} ({w}x{h}px, {len(images_in_slide)} imgs en slide)")
            results.append((slide_idx, name, out_filename))
        except Exception as e:
            print(f"  Slide {slide_idx+1:02d} [{name}]: ERROR al guardar - {e}")
            results.append((slide_idx, name, None))
    
    print(f"\n✅ Proceso completado. Imágenes guardadas en: {OUTPUT_DIR}/")
    print(f"\n{'='*60}")
    print("RESUMEN MAPEO SLIDE -> IMAGEN:")
    print(f"{'='*60}")
    for idx, name, fname in results:
        status = fname if fname else "❌ SIN IMAGEN"
        print(f"  [{idx+1:02d}] {name:<35} -> {status}")

if __name__ == "__main__":
    main()
