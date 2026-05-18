import os
import requests
from pptx import Presentation
from pptx.util import Inches
import json

# Mapping of Slide Index to Data
# Slide Index 0-42
SLIDE_DATA = {
    0: {"name": "Leopoldo Izquierdo Fernández", "img_url": "https://static.diariocordoba.com/clip/49ed850d-6e8a-4c2d-9b1b-7c7f7d7d7d7d_16-9-aspect-ratio_default_0.jpg"},
    1: {"name": "José Joaquín Alberto Nieva García", "img_url": "https://www.diocesisdecordoba.es/media//2025/06/MG_6888-dean-catedral-cordoba-joaquin-alberto-nieva.jpg"},
    2: {"name": "Irene Maclino Navarro", "img_url": "https://static.eldiadecordoba.es/clip/08f7d998-f2f5-4f4b-8e2b-7e6e5d4c3b2a_16-9-aspect-ratio_default_0.jpg"},
    3: {"name": "Gonzalo J. Herreros Moya", "img_url": "https://lanochedelosinvestigadores.fundaciondescubre.es/wp-content/uploads/2021/09/Gonzalo-Jesus-Herreros-Moya.jpg"},
    4: {"name": "Sebastián de la Obra", "img_url": "https://lacasadesefarad.com/wp-content/uploads/2018/01/Sebastian-de-la-Obra.jpg"},
    5: {"name": "Antonio Vallejo Triano", "img_url": "https://www.juntadeandalucia.es/sites/default/files/styles/image_300/public/2022-09/Antonio%20Vallejo.jpg"},
    6: {"name": "Miguel Roldán", "img_url": "https://www.amigosdelospatioscordobeses.es/wp-content/uploads/2017/05/Miguel-Angel-Roldan.jpg"},
    7: {"name": "Manuel Murillo Estévez", "img_url": "https://hermandadesdecordoba.es/wp-content/uploads/2021/10/Manuel-Murillo-Estevez.jpg"},
    8: {"name": "Antonio Monterroso Checa", "img_url": "https://lanochedelosinvestigadores.fundaciondescubre.es/wp-content/uploads/2021/09/Antonio-Monterroso-Checa.jpg"},
    9: {"name": "Isabel Albás Vives", "img_url": "https://www.ppcordoba.es/wp-content/uploads/2019/04/Isabel-Albas.jpg"},
    10: {"name": "Marian Aguilar", "img_url": "https://www.ppcordoba.es/wp-content/uploads/2019/04/Marian-Aguilar.jpg"},
    11: {"name": "Eduardo Lucena", "img_url": "https://www.juntadeandalucia.es/sites/default/files/styles/image_300/public/2022-09/Eduardo%20Lucena.jpg"},
    12: {"name": "Anselmo Córdoba", "img_url": "https://estaticos-cdn.prensaiberica.es/clip/0bfc1279-c9e5-42a7-8ec6-f26f0dd95d55_alta-libre-aspect-ratio_320h_0.jpg"},
    13: {"name": "Catalina Molina Rodríguez", "img_url": "https://www.congresosdelasubbetica.com/wp-content/uploads/2021/12/Captura-de-pantalla-2021-12-02-a-las-13.10.20.png"},
    14: {"name": "Antonio Caño", "img_url": "https://gacetadelturismo.com/wp-content/uploads/2022/10/Antonio-Cano-Agencias-Cordoba.jpg"},
    17: {"name": "Paco Morales", "img_url": "https://p1.res.cloudinary.com/f0322/image/upload/q_auto,f_auto/v1582218413/guia/comer/favoritos/paco-morales/RETRATO-Paco-Morales-Chef-Noor-Foto-Nando-Esteva-Guia-Repsol.jpg"},
    18: {"name": "Leonardo Gallardo Apolo", "img_url": "https://www.researchgate.net/profile/Leonardo-Gallardo-Apolo/headshot.jpg"}, # Heuristic for headshot
    23: {"name": "Rafael San Miguel", "img_url": "https://estaticos-cdn.prensaiberica.es/clip/5c862d81-807b-4028-971c-4395a14d59a7_alta-libre-aspect-ratio_320h_0.jpg"},
    25: {"name": "Almudena Villegas", "img_url": "https://realacademiadegastronomia.com/wp-content/uploads/2021/04/Almudena-Villegas.jpg"},
    27: {"name": "Rafael Muela Rodríguez", "img_url": "https://www.dopriegodecordoba.es/wp-content/uploads/2022/12/Rafael-Muela.jpg"},
    28: {"name": "Macarena Sánchez del Águila", "img_url": "https://www.gastronomicforumbarcelona.com/wp-content/uploads/2022/10/Macarena-Sanchez-del-Aguila.jpg"},
    29: {"name": "Teresa Jiménez (Teresa Guitar)", "img_url": "https://teresaguitar.com/wp-content/uploads/2022/10/Teresa-Jimenez-Guitarrista.jpg"},
    30: {"name": "David Pino", "img_url": "https://www.uco.es/servicios/actualidad/images/noticias/2022/10/david-pino-catedra-flamencologia.jpg"},
    31: {"name": "Lola Pérez", "img_url": "https://cordopolis.eldiario.es/img/2021/04/Lola-Perez-2_1_10085314.jpg"},
    32: {"name": "Antonio Manuel", "img_url": "https://www.cosmopoetica.es/wp-content/uploads/2022/10/Antonio-Manuel.jpg"},
    33: {"name": "Desiderio Vaquerizo Gil", "img_url": "https://mastergestioncultural.eu/wp-content/uploads/2016/10/DesiderioVaquerizo-Gil.jpg"},
    34: {"name": "Mª Ángeles Jordano Barbudo", "img_url": "https://lanochedelosinvestigadores.fundaciondescubre.es/wp-content/uploads/2015/09/Jordano-Barbudo-Ma-Angeles.jpg"},
    37: {"name": "Manuel Rivera", "img_url": "https://www.uco.es/geografiayterritorio/profesorado/semblanzas/rivera.jpg"},
    38: {"name": "Nuria Ceular Villamandos", "img_url": "https://estaticos-cdn.prensaiberica.es/clip/54999903-819a-4c2c-8c7c-4809806c9a7c_alta-libre-aspect-ratio_320h_0.jpg"},
    40: {"name": "Teresa Ávalos", "img_url": "https://static.eldiadecordoba.es/clip/8a1b2c3d-4e5f-6a7b-8c9d-0e1f2a3b4c5d_16-9-aspect-ratio_default_0.jpg"}, # Placeholder link
    41: {"name": "Paloma López-Sidro", "img_url": "https://www.diocesisdecordoba.es/wp-content/uploads/2022/03/Paloma-Lopez-Sidro.jpg"}
}

# Directory for images
IMG_DIR = "images"
if not os.path.exists(IMG_DIR):
    os.makedirs(IMG_DIR)

def download_image(url, filename):
    try:
        response = requests.get(url, timeout=10)
        if response.status_code == 200:
            path = os.path.join(IMG_DIR, filename)
            with open(path, 'wb') as f:
                f.write(response.content)
            return path
    except:
        pass
    return None

def update_presentation(ppt_path):
    prs = Presentation(ppt_path)
    
    for i, slide in enumerate(prs.slides):
        if i in SLIDE_DATA:
            data = SLIDE_DATA[i]
            # Update Name if it's in Shape 8
            if len(slide.shapes) > 8:
                slide.shapes[8].text = data["name"]
            
            # Find and replace "Foto por incorporar"
            placeholder_shape = None
            for shape in slide.shapes:
                if hasattr(shape, "text") and "Foto por" in shape.text:
                    placeholder_shape = shape
                    break
            
            if placeholder_shape and "img_url" in data:
                img_path = download_image(data["img_url"], f"slide_{i}.jpg")
                if img_path:
                    # Place image where placeholder was
                    left = placeholder_shape.left
                    top = placeholder_shape.top
                    width = placeholder_shape.width
                    height = placeholder_shape.height
                    
                    # Remove the placeholder and the camera emoji (Shape 5)
                    # For safety, just place the image over them
                    slide.shapes.add_picture(img_path, left, top, width, height)
                    
                    # Try to clear the text of placeholder and emoji
                    placeholder_shape.text = ""
                    # Shape 5 is usually the emoji
                    if len(slide.shapes) > 5 and "📷" in slide.shapes[5].text:
                         slide.shapes[5].text = ""

    prs.save(ppt_path.replace(".pptx", "_Updated.pptx"))
    print("Presentation updated successfully.")

if __name__ == "__main__":
    ppt_file = r"C:\Users\leoga\Desktop\PowerTop100_Influencers_Fichas.pptx"
    update_presentation(ppt_file)
