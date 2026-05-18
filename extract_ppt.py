from pptx import Presentation
import os
import sys

# Ensure UTF-8 output
if sys.stdout.encoding != 'utf-8':
    sys.stdout.reconfigure(encoding='utf-8')

def extract_content(ppt_path):
    if not os.path.exists(ppt_path):
        print(f"File not found: {ppt_path}")
        return

    prs = Presentation(ppt_path)
    slides_data = []

    for i, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
        slides_data.append({
            "slide_index": i + 1,
            "content": "\n".join(slide_text)
        })

    for data in slides_data:
        print(f"--- Slide {data['slide_index']} ---")
        print(data['content'])
        print("-" * 20)

if __name__ == "__main__":
    ppt_file = "ppt.pptx"
    extract_content(ppt_file)
