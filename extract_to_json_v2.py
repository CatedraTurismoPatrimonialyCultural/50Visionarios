from pptx import Presentation
import os
import sys
import json

# Ensure UTF-8 output
if sys.stdout.encoding != 'utf-8':
    sys.stdout.reconfigure(encoding='utf-8')

def extract_influencers(ppt_path):
    prs = Presentation(ppt_path)
    influencers = []

    for i, slide in enumerate(prs.slides):
        name = None
        entity = None
        description = None
        
        # Based on debug output:
        # Shape 8 is Name
        # Shape 9 is Entity
        # Shape 10 is Description
        try:
            if len(slide.shapes) > 10:
                name = slide.shapes[8].text.strip()
                entity = slide.shapes[9].text.strip()
                description = slide.shapes[10].text.strip()
            else:
                # Fallback search for name/entity in any shape
                # Let's try to find text that isn't the title or placeholder
                potential_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        t = shape.text.strip()
                        if t and "Visionarios" not in t and "Foto por" not in t and "📷" not in t and not t.isdigit():
                            potential_text.append(t)
                if len(potential_text) >= 2:
                    name = potential_text[0]
                    entity = potential_text[1]
                    if len(potential_text) >= 3:
                        description = potential_text[2]
        except:
            pass

        if name and name != "":
            influencers.append({
                "slide_index": i,
                "name": name,
                "entity": entity,
                "description": description
            })

    return influencers

if __name__ == "__main__":
    ppt_file = r"C:\Users\leoga\Desktop\PowerTop100_Influencers_Fichas.pptx"
    data = extract_influencers(ppt_file)
    with open('influencers_full.json', 'w', encoding='utf-8') as f:
        json.dump(data, f, ensure_ascii=False, indent=2)
    print(f"Extracted {len(data)} influencers.")
