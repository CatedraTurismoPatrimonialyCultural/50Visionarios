import os
import re
from pptx import Presentation

def parse_ppt(ppt_path):
    prs = Presentation(ppt_path)
    profiles = []
    
    unsplash_images = [
        "https://images.unsplash.com/photo-1573496359142-b8d87734a5a2?ixlib=rb-4.0.3&auto=format&fit=crop&w=600&q=80",
        "https://images.unsplash.com/photo-1507003211169-0a1dd7228f2d?ixlib=rb-4.0.3&auto=format&fit=crop&w=600&q=80",
        "https://images.unsplash.com/photo-1580489944761-15a19d654956?ixlib=rb-4.0.3&auto=format&fit=crop&w=600&q=80",
        "https://images.unsplash.com/photo-1534528741775-53994a69daeb?ixlib=rb-4.0.3&auto=format&fit=crop&w=600&q=80",
        "https://images.unsplash.com/photo-1560250097-0b93528c311a?ixlib=rb-4.0.3&auto=format&fit=crop&w=600&q=80",
        "https://images.unsplash.com/photo-1573497019940-1c28c88b4f3e?ixlib=rb-4.0.3&auto=format&fit=crop&w=600&q=80",
        "https://images.unsplash.com/photo-1519085360753-af0119f7cbe7?ixlib=rb-4.0.3&auto=format&fit=crop&w=600&q=80",
        "https://images.unsplash.com/photo-1598550874175-4d0ef43ee90d?ixlib=rb-4.0.3&auto=format&fit=crop&w=600&q=80",
        "https://images.unsplash.com/photo-1500648767791-00dcc994a43e?ixlib=rb-4.0.3&auto=format&fit=crop&w=600&q=80",
        "https://images.unsplash.com/photo-1573496799652-408c2ac9fe98?ixlib=rb-4.0.3&auto=format&fit=crop&w=600&q=80"
    ]

    for slide in prs.slides:
        lines = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                lines.extend([l.strip() for l in shape.text.split('\n') if l.strip()])
        
        if not lines: continue
        
        category = "General"
        cat_class = "all"
        name = ""
        role = ""
        bio = ""
        
        for line in lines:
            # Detect category
            if re.match(r'^\d+\.', line):
                category = line
                if "Gestión" in line: cat_class = "gestion"
                elif "Innovación" in line or "Academia" in line: cat_class = "innovacion"
                elif "KOL" in line or "Opinion" in line: cat_class = "kol"
                continue

            if "50 Visionarios" in line or "Turismo Patrimonial" in line and not "Cátedra" in line:
                continue
                
            if not name and len(line) > 4:
                # To avoid picking up random numbers or symbols
                if any(c.isalpha() for c in line):
                    name = line
            elif name and not role and len(line) > 5 and line != name:
                role = line
            elif name and role and line != name and line != role:
                bio += "<p class='mb-3'>" + line + "</p>"
                
        # Clean up possible duplication in bio
        if name and role and bio:
            # Basic deduplication of title in bio
            profiles.append({
                "category": category,
                "cat_class": cat_class,
                "name": name.replace("—", "").strip(),
                "role": role.replace("—", "").strip(),
                "bio": bio,
                "image": unsplash_images[len(profiles) % len(unsplash_images)]
            })

    return profiles

def generate_html(profiles, output_dir="50_visionarios"):
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
        
    html_template_start = """<!DOCTYPE html>
<html lang="es" class="scroll-smooth">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>50 Visionarios del Turismo Patrimonial y Cultural | Córdoba</title>
    <!-- Tailwind CSS -->
    <script src="https://cdn.tailwindcss.com"></script>
    <script>
        tailwind.config = {
            darkMode: 'class',
            theme: {
                extend: {
                    colors: {
                        gold: { 400: '#D4AF37', 500: '#C5A017', 600: '#AA8811' },
                        dark: { 900: '#0a0a0a', 800: '#141414', 700: '#1c1c1c' },
                        olive: '#4A5D23', crimson: '#8B0000', albero: '#FFD700'
                    },
                    fontFamily: {
                        sans: ['Montserrat', 'sans-serif'],
                        serif: ['Playfair Display', 'serif']
                    }
                }
            }
        }
    </script>
    <link href="https://fonts.googleapis.com/css2?family=Montserrat:wght@300;400;600;700&family=Playfair+Display:ital,wght@0,400;0,600;0,700;1,400&display=swap" rel="stylesheet">
    <style>
        body { background-color: #0a0a0a; color: #f3f4f6; }
        .hero-bg {
            background-image: linear-gradient(to bottom, rgba(10, 10, 10, 0.4), rgba(10, 10, 10, 1)), url('https://images.unsplash.com/photo-1549487508-62d10ea3b53f?ixlib=rb-4.0.3&auto=format&fit=crop&w=2000&q=80');
            background-size: cover; background-position: center; background-attachment: fixed;
        }
        .card-image-wrapper { overflow: hidden; }
        .card-image {
            transition: transform 0.6s cubic-bezier(0.25, 0.46, 0.45, 0.94), filter 0.6s ease;
            filter: grayscale(80%) contrast(1.1);
        }
        .group:hover .card-image { transform: scale(1.08); filter: grayscale(0%) contrast(1); }
        ::-webkit-scrollbar { width: 8px; }
        ::-webkit-scrollbar-track { background: #0a0a0a; }
        ::-webkit-scrollbar-thumb { background: #D4AF37; border-radius: 4px; }
        ::-webkit-scrollbar-thumb:hover { background: #C5A017; }
    </style>
</head>
<body class="antialiased selection:bg-gold-500 selection:text-dark-900">

    <nav class="fixed w-full z-50 transition-all duration-300" id="navbar">
        <div class="absolute inset-0 bg-dark-900/80 backdrop-blur-md border-b border-white/5"></div>
        <div class="relative max-w-7xl mx-auto px-6 lg:px-8">
            <div class="flex items-center justify-between h-20">
                <div class="flex-shrink-0 flex items-center gap-4">
                    <span class="font-serif text-2xl font-bold tracking-tight text-white">Cátedra <span class="text-gold-400">Turismo</span></span>
                </div>
                <div class="hidden md:flex items-center space-x-1 font-sans text-sm tracking-widest uppercase">
                    <button class="filter-btn active px-4 py-2 text-gold-400" data-filter="all">Todos</button>
                    <button class="filter-btn px-4 py-2 text-gray-400 hover:text-white" data-filter="gestion">Gestión Patrimonial</button>
                    <button class="filter-btn px-4 py-2 text-gray-400 hover:text-white" data-filter="innovacion">Innovación y Academia</button>
                    <button class="filter-btn px-4 py-2 text-gray-400 hover:text-white" data-filter="kol">KOL / Prescriptores</button>
                </div>
            </div>
        </div>
    </nav>

    <header class="relative hero-bg h-[85vh] flex items-center justify-center text-center">
        <div class="absolute inset-0 bg-gradient-to-t from-dark-900 via-transparent to-transparent"></div>
        <div class="relative z-10 max-w-5xl mx-auto px-6 flex flex-col items-center">
            <p class="text-gold-400 font-sans tracking-[0.2em] uppercase text-sm mb-6">Cátedra de Turismo Patrimonial y Cultural de Córdoba</p>
            <h1 class="font-serif text-5xl md:text-7xl lg:text-8xl text-white font-bold mb-6 leading-tight drop-shadow-lg">
                50 Visionarios <br>
                <span class="text-transparent bg-clip-text bg-gradient-to-r from-gold-400 to-albero font-light italic">del Patrimonio</span>
            </h1>
            <p class="mt-4 text-lg md:text-xl text-gray-300 font-sans max-w-2xl font-light">
                Descubre los perfiles, proyectos e instituciones que están redefiniendo el futuro del turismo cultural y la preservación en Córdoba.
            </p>
        </div>
    </header>

    <main class="max-w-7xl mx-auto px-6 lg:px-8 py-24" id="directorio">
        <div class="flex items-end justify-between mb-16 border-b border-white/10 pb-8">
            <div>
                <h2 class="font-serif text-4xl text-white">Directorio <span class="text-gold-400 italic font-light">de Expertos</span></h2>
                <p class="font-sans text-gray-400 mt-2">Explora los perfiles que impulsan nuestra cultura.</p>
            </div>
        </div>

        <div class="grid grid-cols-1 sm:grid-cols-2 lg:grid-cols-3 xl:grid-cols-4 gap-8" id="grid">
"""

    html_template_end = """
        </div>
    </main>
"""

    footer = """
    <footer class="bg-dark-800 border-t border-white/5 pt-16 pb-8">
        <div class="max-w-7xl mx-auto px-6 lg:px-8">
            <div class="grid grid-cols-1 md:grid-cols-3 gap-12 mb-12">
                <div>
                    <span class="font-serif text-2xl font-bold text-white mb-4 block">Cátedra <span class="text-gold-400">Turismo</span></span>
                    <p class="font-sans text-sm text-gray-400 leading-relaxed max-w-xs">
                        Cátedra de Turismo Patrimonial y Cultural. Impulsando la investigación, gestión y divulgación del legado cordobés.
                    </p>
                </div>
            </div>
            <div class="border-t border-white/5 pt-8 flex flex-col md:flex-row justify-between items-center text-xs font-sans text-gray-500">
                <p>&copy; 2026 Cátedra de Turismo Patrimonial y Cultural de Córdoba.</p>
            </div>
        </div>
    </footer>
"""

    scripts = """
    <script>
        const filterBtns = document.querySelectorAll('.filter-btn');
        const gridItems = document.querySelectorAll('.filter-item');

        filterBtns.forEach(btn => {
            btn.addEventListener('click', () => {
                filterBtns.forEach(b => {
                    b.classList.remove('text-gold-400');
                    b.classList.add('text-gray-400');
                });
                btn.classList.remove('text-gray-400');
                btn.classList.add('text-gold-400');

                const filter = btn.getAttribute('data-filter');

                gridItems.forEach(item => {
                    if (filter === 'all' || item.classList.contains(filter)) {
                        item.style.display = 'flex';
                    } else {
                        item.style.display = 'none';
                    }
                });
            });
        });

        window.addEventListener('scroll', () => {
            const nav = document.getElementById('navbar');
            if (window.scrollY > 50) {
                nav.classList.add('shadow-lg');
            } else {
                nav.classList.remove('shadow-lg');
            }
        });

        function openModal(id) {
            const m = document.getElementById(id);
            if(!m) return; 
            m.classList.remove('hidden'); 
            void m.offsetWidth;
            m.querySelector('.modal-backdrop').classList.replace('opacity-0', 'opacity-100');
            m.querySelector('.modal-content').classList.replace('opacity-0', 'opacity-100');
            m.querySelector('.modal-content').classList.replace('scale-95', 'scale-100');
            document.body.style.overflow = 'hidden';
        }
        function closeModal(id) {
            const m = document.getElementById(id);
            if(!m) return;
            m.querySelector('.modal-backdrop').classList.replace('opacity-100', 'opacity-0');
            m.querySelector('.modal-content').classList.replace('opacity-100', 'opacity-0');
            m.querySelector('.modal-content').classList.replace('scale-100', 'scale-95');
            setTimeout(() => { m.classList.add('hidden'); document.body.style.overflow = ''; }, 300);
        }
        document.addEventListener('keydown', (e) => {
            if (e.key === 'Escape') {
                const openModals = document.querySelectorAll('[id^="modal-"]:not(.hidden)');
                openModals.forEach(modal => closeModal(modal.id));
            }
        });
    </script>
</body>
</html>
"""

    cards_html = ""
    modals_html = ""

    for i, profile in enumerate(profiles):
        idx = i + 1
        cat_short = "Perfil"
        if "innovacion" in profile["cat_class"]: cat_short = "Academia"
        elif "gestion" in profile["cat_class"]: cat_short = "Gestión"
        elif "kol" in profile["cat_class"]: cat_short = "KOL"
        
        cards_html += f"""
            <article class="group cursor-pointer flex flex-col gap-4 filter-item {profile['cat_class']}" onclick="openModal('modal-{idx}')">
                <div class="card-image-wrapper aspect-[3/4] rounded-sm overflow-hidden relative bg-dark-800">
                    <img src="{profile['image']}" alt="Retrato" class="card-image w-full h-full object-cover">
                    <div class="absolute inset-0 bg-gradient-to-t from-dark-900/90 via-dark-900/20 to-transparent opacity-60 group-hover:opacity-80 transition-opacity"></div>
                    <div class="absolute bottom-4 left-4 right-4">
                        <span class="inline-block px-2 py-1 bg-gold-500/20 text-gold-400 text-xs font-sans tracking-wider uppercase backdrop-blur-sm border border-gold-500/30 mb-2">{cat_short}</span>
                    </div>
                </div>
                <div>
                    <h3 class="font-serif text-xl text-white group-hover:text-gold-400 transition-colors">{profile['name']}</h3>
                    <p class="font-sans text-sm text-gray-400 mt-1">{profile['role']}</p>
                </div>
            </article>
        """

        modals_html += f"""
    <div id="modal-{idx}" class="fixed inset-0 z-[100] hidden flex items-center justify-center p-4 sm:p-6" aria-modal="true" role="dialog">
        <div class="absolute inset-0 bg-dark-900/90 backdrop-blur-sm modal-backdrop transition-opacity opacity-0" onclick="closeModal('modal-{idx}')"></div>
        <div class="relative w-full max-w-4xl bg-dark-800 shadow-2xl rounded-sm overflow-hidden flex flex-col md:flex-row modal-content transform scale-95 opacity-0 transition-all duration-300 border border-white/10">
            <button onclick="closeModal('modal-{idx}')" class="absolute top-4 right-4 z-10 text-white/50 hover:text-white bg-dark-900/50 rounded-full p-2 backdrop-blur-md transition-colors">
                ✕
            </button>
            <div class="w-full md:w-2/5 h-64 md:h-auto relative">
                <img src="{profile['image']}" alt="Retrato" class="absolute inset-0 w-full h-full object-cover">
            </div>
            <div class="w-full md:w-3/5 p-8 md:p-12 overflow-y-auto max-h-[80vh]">
                <span class="text-gold-400 font-sans tracking-widest uppercase text-xs font-semibold mb-2 block">{profile['category']}</span>
                <h2 class="font-serif text-3xl md:text-4xl text-white mb-2">{profile['name']}</h2>
                <p class="font-sans text-gray-400 mb-8 pb-6 border-b border-white/10">{profile['role']}</p>
                <div class="prose prose-invert prose-p:font-sans prose-p:font-light prose-p:leading-relaxed prose-p:text-gray-300">
                    {profile['bio']}
                </div>
            </div>
        </div>
    </div>
        """

    full_html = html_template_start + cards_html + html_template_end + modals_html + footer + scripts

    with open(os.path.join(output_dir, "index.html"), "w", encoding="utf-8") as f:
        f.write(full_html)
    print(f"Generated {len(profiles)} profiles in {output_dir}/index.html")

if __name__ == '__main__':
    profiles = parse_ppt('ppt.pptx')
    generate_html(profiles, '50_visionarios')
