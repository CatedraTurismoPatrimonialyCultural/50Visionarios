from pptx import Presentation
import os
import sys

# Ensure UTF-8 output
if sys.stdout.encoding != 'utf-8':
    sys.stdout.reconfigure(encoding='utf-8')

def debug_slides(ppt_path):
    prs = Presentation(ppt_path)
    for i, slide in enumerate(prs.slides):
        if i >= 3: break # Check first 3
        print(f"--- Slide {i+1} ---")
        for j, shape in enumerate(slide.shapes):
            if hasattr(shape, "text"):
                print(f"Shape {j} Text: {shape.text!r}")
        print("-" * 20)

if __name__ == "__main__":
    ppt_file = r"C:\Users\leoga\Desktop\PowerTop100_Influencers_Fichas.pptx"
    debug_slides(ppt_file)
