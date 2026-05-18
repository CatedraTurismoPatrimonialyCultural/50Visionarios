from pptx import Presentation
import os
import sys
import json

# Ensure UTF-8 output
if sys.stdout.encoding != 'utf-8':
    sys.stdout.reconfigure(encoding='utf-8')

def extract_influencers(ppt_path):
    prs = Presentation(ppt_path)
    influencers = []

    for i, slide in enumerate(prs.slides):
        lines = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                lines.extend(shape.text.split('\n'))
        
        # Heuristic to find name and entity
        # Looking for the block after "Foto por incorporar"
        name = None
        entity = None
        description = None
        
        try:
            # Simple heuristic based on the structure observed
            # Name is usually the first line after "Foto por incorporar"
            found_photo = False
            for line in lines:
                l = line.strip()
                if "Foto por incorporar" in l:
                    found_photo = True
                    continue
                if found_photo and l and not name:
                    name = l
                    continue
                if name and not entity and l:
                    entity = l
                    continue
                if entity and not description and l:
                    description = l
                    continue
        except:
            pass

        if name:
            influencers.append({
                "slide_index": i,
                "name": name,
                "entity": entity,
                "description": description
            })

    return influencers

if __name__ == "__main__":
    ppt_file = r"C:\Users\leoga\Desktop\PowerTop100_Influencers_Fichas.pptx"
    data = extract_influencers(ppt_file)
    with open('influencers.json', 'w', encoding='utf-8') as f:
        json.dump(data, f, ensure_ascii=False, indent=2)
    print(f"Extracted {len(data)} influencers.")
